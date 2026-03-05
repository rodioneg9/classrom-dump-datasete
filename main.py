import os
import json
import re
import time
import pdfplumber
from docx import Document
from pptx import Presentation

from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build


COURSE_ID = "ВСТАВЬТЕ_ID_КУРСА"

DATASET_DIR = "dataset"
DATASET_JSONL = "dataset_llm.jsonl"


def safe_name(name):

    if not name:
        return "no_name"

    name = name.replace("\n", " ").replace("\r", " ")

    name = re.sub(r'[\\/*?:"<>|]', "", name)

    name = re.sub(r"\s+", " ", name)

    name = name.strip()

    name = name.rstrip(". ")

    # ограничиваем длину имени
    name = name[:60]

    if not name:
        name = "no_name"

    return name


def retry(func, retries=5):

    for i in range(retries):

        try:
            return func()

        except Exception as e:

            print("Retry", i + 1, "error:", e)
            time.sleep(5)

    raise Exception("Max retries exceeded")


def extract_pdf(path):

    text = ""

    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"

    except:
        pass

    return text


def extract_docx(path):

    text = ""

    try:
        doc = Document(path)

        for p in doc.paragraphs:
            text += p.text + "\n"

    except:
        pass

    return text


def extract_pptx(path):

    text = ""

    try:
        prs = Presentation(path)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    except:
        pass

    return text


def extract_text(path):

    ext = path.lower().split(".")[-1]

    if ext == "pdf":
        return extract_pdf(path)

    if ext == "docx":
        return extract_docx(path)

    if ext == "pptx":
        return extract_pptx(path)

    return ""


def find_last_downloaded(course_dir):

    max_id = 0

    if not os.path.exists(course_dir):
        return 0

    for folder in os.listdir(course_dir):

        m = re.match(r"(\d+)_", folder)

        if m:

            n = int(m.group(1))

            if n > max_id:
                max_id = n

    return max_id


creds = Credentials.from_authorized_user_file("token.json")

classroom = build("classroom", "v1", credentials=creds)
drive = build("drive", "v3", credentials=creds)

os.makedirs(DATASET_DIR, exist_ok=True)

course = retry(lambda: classroom.courses().get(id=COURSE_ID).execute())

course_name = safe_name(course["name"])

course_dir = os.path.join(DATASET_DIR, course_name)

os.makedirs(course_dir, exist_ok=True)

print("Course:", course_name)


assignment_index = find_last_downloaded(course_dir) + 1

print("Resume from assignment:", assignment_index)


dataset_records = []

page_token = None
global_index = 1


while True:

    response = retry(lambda: classroom.courses().courseWork().list(
        courseId=COURSE_ID,
        pageToken=page_token
    ).execute())

    coursework = response.get("courseWork", [])

    for work in coursework:

        if global_index < assignment_index:

            global_index += 1
            continue


        title = safe_name(work.get("title", "no_title"))

        folder_name = f"{global_index:04d}_{title}"

        work_dir = os.path.join(course_dir, folder_name)

        # защита от пробелов и точек в конце
        work_dir = work_dir.rstrip(" .")

        # защита от слишком длинных путей Windows
        if len(work_dir) > 200:
            work_dir = work_dir[:200]

        if os.path.exists(work_dir):

            print("skip existing assignment:", title)

            global_index += 1
            continue


        os.makedirs(work_dir, exist_ok=True)

        print(f"{global_index:04d}", title)

        description = work.get("description", "")

        task_path = os.path.join(work_dir, "task.json")

        with open(task_path, "w", encoding="utf8") as f:
            json.dump(work, f, indent=2, ensure_ascii=False)

        combined_text = f"{title}\n{description}\n"


        if "materials" in work:

            for m in work["materials"]:

                if "driveFile" not in m:
                    continue

                file_id = m["driveFile"]["driveFile"]["id"]

                try:

                    meta = retry(lambda: drive.files().get(
                        fileId=file_id,
                        fields="name,mimeType"
                    ).execute())

                    file_name = safe_name(meta["name"])

                    mime = meta["mimeType"]

                    path = os.path.join(work_dir, file_name)

                    print("   file:", file_name)


                    if mime == "application/vnd.google-apps.document":

                        path += ".pdf"

                        if os.path.exists(path):
                            continue

                        request = drive.files().export(
                            fileId=file_id,
                            mimeType="application/pdf"
                        )


                    elif mime == "application/vnd.google-apps.spreadsheet":

                        path += ".xlsx"

                        if os.path.exists(path):
                            continue

                        request = drive.files().export(
                            fileId=file_id,
                            mimeType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        )


                    elif mime == "application/vnd.google-apps.presentation":

                        path += ".pptx"

                        if os.path.exists(path):
                            continue

                        request = drive.files().export(
                            fileId=file_id,
                            mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )


                    else:

                        if os.path.exists(path):
                            continue

                        request = drive.files().get_media(fileId=file_id)


                    data = retry(lambda: request.execute())

                    with open(path, "wb") as f:
                        f.write(data)


                    extracted = extract_text(path)

                    if extracted:
                        combined_text += "\n" + extracted


                except Exception as e:

                    print("download error:", e)


        dataset_records.append({
            "course": course_name,
            "assignment": title,
            "text": combined_text.strip()
        })


        global_index += 1


    page_token = response.get("nextPageToken")

    if not page_token:
        break


with open(DATASET_JSONL, "a", encoding="utf8") as out:

    for rec in dataset_records:
        out.write(json.dumps(rec, ensure_ascii=False) + "\n")


print("\nDownload finished or resumed successfully.")
